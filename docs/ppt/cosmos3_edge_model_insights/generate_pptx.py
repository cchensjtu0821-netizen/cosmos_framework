#!/usr/bin/env python3
"""Generate the editable Cosmos3-Edge analysis slide and its PNG preview.

Dependencies are intentionally not vendored. Run with python-pptx and Pillow on
PYTHONPATH; the generated PPTX uses native PowerPoint shapes and connectors.
"""

from __future__ import annotations

import math
import os
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
PPTX_OUT = ROOT / "cosmos3_edge_model_insights.pptx"
PNG_OUT_PAGE1 = ROOT / "cosmos3_edge_model_insights_page1.png"
PNG_OUT_PAGE2 = ROOT / "cosmos3_edge_model_insights_page2.png"
WIDTH, HEIGHT = 1920, 1080
PX_PER_INCH = 144
FONT_FILE = Path(os.environ.get("COSMOS_PPT_FONT", "/tmp/NotoSansCJKsc-Regular.otf"))


COLORS = {
    "canvas": "F5F7FB",
    "panel": "FFFFFF",
    "ink": "14213D",
    "muted": "5E6B82",
    "hair": "D9E2EF",
    "grid": "EAF0F7",
    "navy": "183B66",
    "blue": "1677FF",
    "cyan": "09A9C8",
    "green": "20A464",
    "amber": "E69A24",
    "red": "D94A5A",
    "soft_blue": "EAF3FF",
    "soft_cyan": "E7F8FB",
    "soft_green": "EAF8F1",
    "soft_amber": "FFF5E4",
    "soft_red": "FDECEF",
    "soft_navy": "EDF2F8",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def inch(px: float):
    return Inches(px / PX_PER_INCH)


def ppt_font_size(px: float):
    return Pt(px / 2)


def font(px: int, *, bold: bool = False):
    path = FONT_FILE if FONT_FILE.exists() else Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
    return ImageFont.truetype(str(path), px, layout_engine=ImageFont.Layout.RAQM)


def hex_tuple(hex_color: str):
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


class SlideCanvas:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = inch(WIDTH)
        self.prs.slide_height = inch(HEIGHT)
        self.images = []
        self.new_slide()

    def new_slide(self):
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.image = Image.new("RGB", (WIDTH, HEIGHT), hex_tuple(COLORS["canvas"]))
        self.draw = ImageDraw.Draw(self.image)
        self.images.append(self.image)
        self._set_background()

    def _set_background(self):
        bg = self.slide.background.fill
        bg.solid()
        bg.fore_color.rgb = rgb(COLORS["canvas"])
        for x in range(0, WIDTH, 48):
            self.line(x, 0, x, HEIGHT, COLORS["grid"], 1)
        for y in range(0, HEIGHT, 48):
            self.line(0, y, WIDTH, y, COLORS["grid"], 1)

    def rect(self, x, y, w, h, fill, stroke="D9E2EF", radius=14, width=2):
        shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = self.slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(stroke)
        shape.line.width = Pt(max(0.5, width / 2))
        self.draw.rounded_rectangle((x, y, x + w, y + h), radius=radius, fill=hex_tuple(fill), outline=hex_tuple(stroke), width=width)
        return shape

    def line(self, x1, y1, x2, y2, color, width=2, dash=False):
        connector = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
        connector.line.color.rgb = rgb(color)
        connector.line.width = Pt(max(0.5, width / 2))
        if dash:
            connector.line.dash_style = 4
            self.draw.line((x1, y1, x2, y2), fill=hex_tuple(color), width=width)
        else:
            self.draw.line((x1, y1, x2, y2), fill=hex_tuple(color), width=width)
        return connector

    def arrow(self, x1, y1, x2, y2, color="1677FF", width=4, dash=False):
        self.line(x1, y1, x2, y2, color, width, dash)
        angle = math.atan2(y2 - y1, x2 - x1)
        size = 12
        left = (x2 - size * math.cos(angle - math.pi / 6), y2 - size * math.sin(angle - math.pi / 6))
        right = (x2 - size * math.cos(angle + math.pi / 6), y2 - size * math.sin(angle + math.pi / 6))
        tri = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, inch(x2 - 9), inch(y2 - 9), inch(18), inch(18))
        tri.rotation = math.degrees(angle) + 90
        tri.fill.solid()
        tri.fill.fore_color.rgb = rgb(color)
        tri.line.fill.background()
        self.draw.polygon([(x2, y2), left, right], fill=hex_tuple(color))

    def path_arrow(self, points, color="1677FF", width=4, dash=False):
        """Draw an editable orthogonal/segmented connector with one arrowhead."""
        for (x1, y1), (x2, y2) in zip(points, points[1:]):
            self.line(x1, y1, x2, y2, color, width, dash)
        (x1, y1), (x2, y2) = points[-2], points[-1]
        angle = math.atan2(y2 - y1, x2 - x1)
        size = 12
        left = (x2 - size * math.cos(angle - math.pi / 6), y2 - size * math.sin(angle - math.pi / 6))
        right = (x2 - size * math.cos(angle + math.pi / 6), y2 - size * math.sin(angle + math.pi / 6))
        tri = self.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, inch(x2 - 9), inch(y2 - 9), inch(18), inch(18))
        tri.rotation = math.degrees(angle) + 90
        tri.fill.solid()
        tri.fill.fore_color.rgb = rgb(color)
        tri.line.fill.background()
        self.draw.polygon([(x2, y2), left, right], fill=hex_tuple(color))

    def text(self, x, y, w, h, value, size=24, color="14213D", bold=False, align="left", valign="top", margin=0, line_spacing=1.0):
        box = self.slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = inch(margin)
        frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
        p = frame.paragraphs[0]
        p.text = value
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        p.font.name = "Microsoft YaHei"
        p.font.size = ppt_font_size(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing

        pil_font = font(size, bold=bold)
        anchor = {"left": "la", "center": "ma", "right": "ra"}[align]
        tx = {"left": x + margin, "center": x + w / 2, "right": x + w - margin}[align]
        bbox = self.draw.multiline_textbbox((0, 0), value, font=pil_font, spacing=max(2, int(size * 0.25)), align=align)
        text_h = bbox[3] - bbox[1]
        ty = {"top": y + margin, "middle": y + (h - text_h) / 2, "bottom": y + h - text_h - margin}[valign]
        self.draw.multiline_text((tx, ty), value, font=pil_font, fill=hex_tuple(color), anchor=anchor, spacing=max(2, int(size * 0.25)), align=align)
        return box

    def badge(self, x, y, w, label, fill, fg="FFFFFF"):
        self.rect(x, y, w, 34, fill, fill, 17, 1)
        self.text(x, y, w, 34, label, 17, fg, True, "center", "middle")

    def node(self, x, y, w, h, title, detail, fill="FFFFFF", stroke="9EB8D7", accent="1677FF", flag=None, title_size=20):
        self.rect(x, y, w, h, fill, stroke, 12, 2)
        self.rect(x, y, 7, h, accent, accent, 3, 1)
        self.text(x + 18, y + 11, w - 30, 26, title, title_size, COLORS["ink"], True)
        self.text(x + 18, y + 41, w - 30, h - 50, detail, 16, COLORS["muted"], False, "left", "top", 0, 0.9)
        if flag:
            self.badge(x + w - 54, y + 9, 42, flag, accent)

    def set_notes(self, value):
        notes = self.slide.notes_slide.notes_text_frame
        notes.text = value

    def save(self):
        props = self.prs.core_properties
        props.title = "NVIDIA Cosmos3-Edge 模型洞察与竞品分析"
        props.subject = "模型分析情况与 H100 实测情况"
        props.author = "Cosmos Framework"
        self.prs.save(PPTX_OUT)
        preview_paths = [PNG_OUT_PAGE1, PNG_OUT_PAGE2]
        for image, path in zip(self.images, preview_paths):
            image.save(path, optimize=True)


def build_h100_slide(c):
    c.new_slide()

    # Header
    c.badge(56, 40, 86, "BENCH 02", COLORS["navy"])
    c.text(160, 34, 1180, 60, "H100 实测情况", 46, COLORS["ink"], True)
    c.text(160, 95, 1180, 32, "RoboLab 远端闭环 · 模块级 profiling · 性能瓶颈定位", 21, COLORS["muted"])
    c.text(1460, 45, 404, 32, "EDGE POLICY · DROID", 18, COLORS["blue"], True, "right")
    c.text(1460, 79, 404, 28, "H100 server / RTX 4090 client", 16, COLORS["muted"], False, "right")
    c.line(56, 140, 1864, 140, COLORS["hair"], 2)

    # 01 Test environment and method
    ex, ey, ew, eh = 56, 166, 1212, 332
    c.rect(ex, ey, ew, eh, COLORS["panel"], COLORS["hair"], 18, 2)
    c.badge(ex + 24, ey + 20, 44, "01", COLORS["blue"])
    c.text(ex + 82, ey + 15, 360, 38, "测试环境与方法", 27, COLORS["ink"], True)
    method_badges = [
        (ex + 544, "4-step · CFG=3", COLORS["blue"]),
        (ex + 708, "稳定样本 n=2", COLORS["green"]),
        (ex + 872, "decode video = false", COLORS["amber"]),
    ]
    for xx, label, fill in method_badges:
        c.badge(xx, ey + 19, 146 if "decode" not in label else 196, label, fill)

    # RTX 4090 RoboLab client
    c.rect(88, 246, 286, 162, COLORS["soft_blue"], "8FB9ED", 14, 3)
    c.badge(104, 260, 98, "RTX 4090", COLORS["blue"])
    c.text(104, 303, 254, 28, "RoboLab Client", 22, COLORS["ink"], True, "center")
    c.text(104, 338, 254, 58, "仿真 / 相机观测 / prompt\n执行 action → 下一帧", 15, COLORS["muted"], False, "center", "middle")

    # SSH tunnel / WebSocket bridge
    c.arrow(378, 312, 438, 312, COLORS["blue"], 4)
    c.rect(442, 270, 208, 88, COLORS["soft_amber"], "E8BE6A", 12, 2)
    c.text(454, 278, 184, 26, "SSH 端口转发", 18, COLORS["ink"], True, "center")
    c.text(454, 311, 184, 36, "WebSocket request / response\n端口不写死", 13, COLORS["muted"], False, "center")
    c.arrow(654, 312, 704, 312, COLORS["amber"], 4)
    c.text(378, 287, 60, 18, "OBS", 12, COLORS["blue"], True, "center")

    # H100 policy server
    c.rect(708, 232, 516, 184, COLORS["soft_green"], "90CFAD", 14, 3)
    c.badge(724, 246, 86, "H100", COLORS["green"])
    c.text(828, 246, 370, 30, "Cosmos3-Edge Policy Server", 22, COLORS["ink"], True)
    server_stages = [
        (728, 294, 132, "PREPARE", "VAE / pack"),
        (872, 294, 164, "SAMPLER", "UniPC ×4 · CFG"),
        (1048, 294, 156, "OUTPUT", "action [32,8]"),
    ]
    for xx, yy, ww, title, detail in server_stages:
        c.rect(xx, yy, ww, 70, "FFFFFF", "B9DCC9", 9, 2)
        c.text(xx + 8, yy + 9, ww - 16, 20, title, 14, COLORS["green"], True, "center")
        c.text(xx + 8, yy + 35, ww - 16, 20, detail, 13, COLORS["muted"], False, "center")

    # Returned action closes the robot-control loop.
    c.path_arrow([(1112, 420), (1112, 460), (232, 460), (232, 412)], COLORS["green"], 4)
    c.rect(510, 445, 322, 30, COLORS["green"], COLORS["green"], 15, 1)
    c.text(510, 445, 322, 30, "ACTION [32,8] · CLOSED LOOP", 13, "FFFFFF", True, "center", "middle")
    c.text(916, 429, 186, 20, "WebSocket response", 12, COLORS["green"], True, "right")

    # 02 Model-effect placeholders owned by the user.
    mx, my, mw, mh = 1290, 166, 574, 332
    c.rect(mx, my, mw, mh, COLORS["panel"], COLORS["hair"], 18, 2)
    c.badge(mx + 24, my + 20, 44, "02", COLORS["amber"])
    c.text(mx + 82, my + 15, 360, 38, "模型效果", 27, COLORS["ink"], True)
    c.badge(mx + 438, my + 19, 112, "待补素材", COLORS["amber"])
    effect_cards = [
        (mx + 24, "生成适配效果", "放生成帧 / 对齐结果"),
        (mx + 292, "机器人操作效果", "放实机图 / 视频二维码"),
    ]
    for xx, title, detail in effect_cards:
        c.rect(xx, my + 78, 258, 220, "F8FAFD", "BFCBDD", 12, 2)
        c.text(xx, my + 102, 258, 58, "+", 48, "AAB7C8", False, "center", "middle")
        c.text(xx + 14, my + 172, 230, 28, title, 18, COLORS["ink"], True, "center")
        c.text(xx + 14, my + 207, 230, 48, detail, 14, COLORS["muted"], False, "center", "middle")
        c.badge(xx + 58, my + 264, 142, "USER CONTENT", COLORS["navy"])

    # 03 Performance data and bottleneck analysis
    px, py, pw, ph = 56, 520, 1808, 502
    c.rect(px, py, pw, ph, COLORS["panel"], COLORS["hair"], 18, 2)
    c.badge(px + 24, py + 20, 44, "03", COLORS["green"])
    c.text(px + 82, py + 15, 600, 38, "性能数据与瓶颈分析", 27, COLORS["ink"], True)
    c.text(px + 860, py + 22, 900, 26, "纯模型推理口径｜不含 WebSocket / SSH / 服务包装", 15, COLORS["muted"], False, "right")

    kpis = [
        ("312.75 ms", "generator_total · 模型推理", COLORS["red"], COLORS["soft_red"]),
        ("3.197 inf/s", "单实例串行推理吞吐", COLORS["blue"], COLORS["soft_blue"]),
        ("231.83 ms", "sampler · 74.13% 推理", COLORS["amber"], COLORS["soft_amber"]),
        ("117.087 TF", "理论计算量 / 次推理", COLORS["green"], COLORS["soft_green"]),
        ("8,438.66 MiB", "peak alloc · reserved 8,644", COLORS["navy"], COLORS["soft_navy"]),
    ]
    kx, ky, kw, kg = px + 28, py + 72, 334, 14
    for i, (metric, label, accent, fill) in enumerate(kpis):
        xx = kx + i * (kw + kg)
        c.rect(xx, ky, kw, 98, fill, accent, 12, 2)
        c.rect(xx, ky, 8, 98, accent, accent, 4, 1)
        c.text(xx + 22, ky + 11, kw - 38, 40, metric, 28, accent, True)
        c.text(xx + 22, ky + 54, kw - 38, 28, label, 14, COLORS["muted"])

    def draw_table(x, y, title, columns, headers, rows, header_fill):
        c.text(x, y, sum(columns), 26, title, 17, COLORS["ink"], True)
        xx, yy = x, y + 30
        for label, ww in zip(headers, columns):
            c.rect(xx, yy, ww, 31, header_fill, header_fill, 0, 1)
            c.text(xx + 4, yy, ww - 8, 31, label, 12, "FFFFFF", True, "center", "middle")
            xx += ww
        yy += 31
        for ri, row in enumerate(rows):
            xx = x
            for ci, (value, ww) in enumerate(zip(row, columns)):
                fill = "F6F9FD" if ri % 2 == 0 else "FFFFFF"
                if ci == 0:
                    fill = COLORS["soft_navy"]
                c.rect(xx, yy, ww, 26, fill, COLORS["hair"], 0, 1)
                c.text(xx + 4, yy, ww - 8, 26, value, 11, COLORS["ink"] if ci == 0 else COLORS["muted"], ci == 0, "center", "middle")
                xx += ww
            yy += 26

    table_y = py + 188
    latency_rows = [
        ("generator_total", "312.75", "100% 推理", "prepare + sampler 总区间"),
        ("prepare_data_total", "75.90", "24.27% 推理", "condition / VAE / initial pack"),
        ("vae_encode", "9.92", "3.17% 推理", "视觉条件编码，非主瓶颈"),
        ("sampler_total", "231.83", "74.13% 推理", "4-step · 4C + 4U"),
        ("network_forward", "222.635", "71.19% 推理", "占 sampler 96.03%"),
        ("mot_joint_forward", "178.407", "57.05% 推理", "占 network 80.13%"),
        ("encode_vision", "7.478", "2.39% 推理", "视觉 latent 输入投影"),
        ("build_attention", "6.399", "2.05% 推理", "mask / layout / 索引构造"),
    ]
    draw_table(
        px + 28,
        table_y,
        "A｜稳定态模型推理时间（ms）",
        [220, 130, 170, 330],
        ["模块", "平均耗时", "占模型推理", "结论"],
        latency_rows,
        COLORS["navy"],
    )

    flops_rows = [
        ("generator_total", "117.087", "100.00%", "374.4 TF/s · 推理总区间"),
        ("VAE encode", "1.074", "0.92%", "108.3 TF/s"),
        ("sampler / network", "116.012", "99.08%", "500.4 TF/s · sampler"),
        ("MoT joint forward", "115.973", "99.05%", "650.0 TF/s · MoT"),
        ("encode_vision", "0.01933", "0.0165%", "低 FLOPs / 7.478 ms"),
        ("vision_head", "0.01925", "0.0164%", "低 FLOPs / 0.686 ms"),
        ("encode_action", "0.000145", "~0.00012%", "低 FLOPs / 4.813 ms"),
        ("action_head", "0.000069", "~0.00006%", "低 FLOPs / 0.786 ms"),
    ]
    draw_table(
        px + 906,
        table_y,
        "B｜理论计算量与折算速率",
        [220, 160, 160, 310],
        ["模块", "TFLOPs / 次推理", "占推理 FLOPs", "稳定态折算速率 / 观察"],
        flops_rows,
        COLORS["green"],
    )

    c.text(px + 28, py + 462, 1710, 20, "结论｜MoT 承担 99.05% 理论 FLOPs、占模型推理 57.05% 时间；Sampler 占 74.13%，CFG 每次推理执行 4C + 4U。", 13, COLORS["red"], True)

    c.text(56, 1037, 1660, 22, "Source: RoboLab Edge Policy Server 推理性能分析报告 · stable samples 4–5 · 2026-07-27", 13, COLORS["muted"])
    c.text(1720, 1037, 144, 22, "2026.08", 13, COLORS["muted"], True, "right")

    c.set_notes(
        "讲解要点：RoboLab 客户端运行在 RTX 4090，观测通过 SSH 端口转发后的 WebSocket 发往 H100 上的 "
        "Cosmos3-Edge Policy Server，返回 [32,8] action 后闭环执行；该拓扑只解释测试环境，不进入性能口径。"
        "本页从 generator_total 开始统计纯模型推理，排除 WebSocket、SSH、build_sample 和服务包装。稳定模型推理为 312.75 ms，"
        "对应单实例串行 3.197 inference/s。4-step、guidance=3 每次推理执行 4 次 conditional 和 4 次 unconditional。"
        "network_forward 的 222.635 ms 是 8 次调用的累计值，sampler 总计 231.83 ms。理论计算量 117.087 TFLOPs/次推理，其中 MoT joint forward 为 115.973 TFLOPs，"
        "是首要 GPU 瓶颈。稳定峰值显存为 8438.66 MiB allocated、8644 MiB reserved，request 2–5 无增长。"
        "模型效果区域由用户补入素材。"
    )


def build():
    c = SlideCanvas()

    # Header
    c.badge(56, 40, 86, "MODEL 01", COLORS["navy"])
    c.text(160, 34, 1180, 60, "NVIDIA Cosmos3-Edge｜模型洞察与竞品分析", 46, COLORS["ink"], True)
    c.text(160, 95, 1180, 32, "模型分析情况：单次推理架构 · 参数/计算结构 · 具身关键创新", 21, COLORS["muted"])
    c.text(1460, 45, 404, 32, "EDGE POLICY · DROID", 18, COLORS["blue"], True, "right")
    c.text(1460, 79, 404, 28, "32-step observed runtime", 16, COLORS["muted"], False, "right")
    c.line(56, 140, 1864, 140, COLORS["hair"], 2)

    # Left architecture panel
    lx, ly, lw, lh = 56, 166, 1128, 856
    c.rect(lx, ly, lw, lh, COLORS["panel"], COLORS["hair"], 18, 2)
    c.badge(lx + 26, ly + 22, 44, "01", COLORS["blue"])
    c.text(lx + 84, ly + 17, 760, 38, "Edge 单次推理框图", 28, COLORS["ink"], True)
    c.text(lx + 810, ly + 23, 286, 28, "去除 WebSocket I/O", 16, COLORS["muted"], False, "right")

    # Step 1: collapse all service-side preprocessing into one audience-level block.
    c.badge(84, 242, 88, "STEP 1", COLORS["navy"])
    c.text(184, 242, 190, 30, "输入汇聚", 20, COLORS["ink"], True)
    c.rect(84, 282, 216, 258, COLORS["soft_navy"], "AFC2D9", 14, 2)
    c.text(102, 296, 180, 30, "Input & Batch Prep", 20, COLORS["ink"], True, "center")
    c.text(102, 328, 180, 34, "build_sample / transform / batch", 12, COLORS["muted"], False, "center")
    input_lanes = [
        (374, "VIDEO", "[1,3,33,544,736]", COLORS["cyan"], COLORS["soft_cyan"]),
        (436, "ACTION / STATE", "batch [1,33,64] · sample [33,64]", COLORS["amber"], COLORS["soft_amber"]),
        (498, "PROMPT", "cond / uncond text", COLORS["blue"], COLORS["soft_blue"]),
    ]
    for yy, label, shape, accent, fill in input_lanes:
        c.rect(102, yy, 180, 50, fill, accent, 9, 2)
        c.text(112, yy + 4, 160, 19, label, 13, accent, True)
        c.text(112, yy + 23, 160, 20, shape, 12, COLORS["muted"])

    # Step 2: preparation is a coherent stage with three modality lanes and one merge.
    c.rect(320, 236, 808, 334, "F9FBFE", "8FB0D3", 15, 3)
    c.badge(338, 252, 88, "STEP 2", COLORS["blue"])
    c.badge(438, 252, 176, "prepare_data_total", COLORS["navy"])
    c.text(630, 254, 466, 26, "模态编码 → 顺序打包 → 联合扩散初值", 17, COLORS["muted"], True)

    c.node(348, 314, 230, 108, "VAE Encoder", "仅 frame 0 前缀编码\nraw [1,48,9,34,46]\ncrop → [1,48,9,33,40]", COLORS["soft_cyan"], "8CCED9", COLORS["cyan"], "V")
    c.node(348, 454, 230, 82, "Tokenize Text", "cond 105→107 · uncond 17→19", COLORS["soft_blue"], "8FB9ED", COLORS["blue"], "T")

    # Explicit input data routes and shapes.
    c.arrow(302, 399, 344, 365, COLORS["cyan"], 4)
    c.path_arrow([(302, 461), (322, 461), (322, 548), (642, 548), (642, 496), (658, 496)], COLORS["amber"], 4)
    c.arrow(302, 523, 344, 495, COLORS["blue"], 4)
    c.text(414, 532, 206, 18, "ACTION LATENT [33,64]", 12, COLORS["amber"], True, "center")

    # Merge initial_pack and initialize-noise into one logical stage.
    c.rect(662, 304, 438, 232, COLORS["soft_amber"], "E8BE6A", 14, 3)
    c.badge(680, 320, 192, "PACK + NOISE INIT", COLORS["amber"])
    c.text(884, 321, 198, 28, "严格 token 顺序", 16, COLORS["ink"], True, "right")
    c.arrow(582, 366, 658, 376, COLORS["cyan"], 4)
    c.text(586, 342, 70, 20, "V-LATENT", 11, COLORS["cyan"], True, "center")
    c.arrow(582, 496, 658, 454, COLORS["blue"], 4)
    c.text(584, 476, 72, 20, "TEXT IDS", 11, COLORS["blue"], True, "center")

    # The standard Policy packer is text, then vision, then action.
    token_cells = [
        (682, 370, 112, "① TEXT", "C107 / U19", COLORS["blue"], COLORS["soft_blue"]),
        (798, 370, 144, "② VISION", "3060 patches", COLORS["cyan"], COLORS["soft_cyan"]),
        (946, 370, 132, "③ ACTION", "33 tokens", COLORS["green"], COLORS["soft_green"]),
    ]
    for xx, yy, ww, label, detail, accent, fill in token_cells:
        c.rect(xx, yy, ww, 60, fill, accent, 8, 2)
        c.text(xx, yy + 7, ww, 20, label, 13, accent, True, "center")
        c.text(xx, yy + 31, ww, 18, detail, 12, COLORS["muted"], False, "center")
    c.arrow(794, 400, 796, 400, COLORS["blue"], 2)
    c.arrow(942, 400, 944, 400, COLORS["cyan"], 2)
    c.text(682, 441, 396, 20, "packed hidden: cond [3200,2048] · uncond [3112,2048]", 13, COLORS["navy"], True, "center")
    c.rect(682, 470, 396, 48, COLORS["soft_red"], "E9A9B1", 8, 2)
    c.text(692, 474, 376, 18, "JOINT NOISE / LATENT", 12, COLORS["red"], True, "center")
    c.text(692, 494, 376, 18, "noise_v [1,48,9,33,40]  +  noise_a [33,64]  → flat [572352]", 12, COLORS["muted"], False, "center")

    # One unmistakable hand-off from preparation into the sampler.
    c.path_arrow([(880, 540), (880, 588), (100, 588), (100, 664)], COLORS["navy"], 5)
    c.rect(408, 573, 354, 30, COLORS["navy"], COLORS["navy"], 15, 1)
    c.text(408, 573, 354, 30, "PACKED CONDITIONS + JOINT LATENT", 13, "FFFFFF", True, "center", "middle")

    # Step 3: sampler_total is the main denoising system boundary.
    sy, sh = 610, 258
    c.rect(84, sy, 1044, sh, "F9FBFE", "8FB0D3", 15, 3)
    c.badge(112, sy + 15, 88, "STEP 3", COLORS["green"])
    c.badge(212, sy + 15, 142, "sampler_total", COLORS["navy"])
    c.text(372, sy + 18, 520, 26, "双分支 MoT 去噪 → CFG 合流 → joint latent 更新", 17, COLORS["muted"], True)
    c.badge(946, sy + 14, 158, "UniPC × 4 steps", COLORS["green"])

    c.node(116, 676, 314, 72, "Conditional MoT forward", "C [3200,2048] → v_c{video, action}", COLORS["soft_blue"], "8FB9ED", COLORS["blue"], "C", 17)
    c.node(116, 770, 314, 70, "Unconditional MoT forward", "U [3112,2048] → v_u · guidance≠1", COLORS["soft_navy"], "AFC2D9", COLORS["navy"], "U", 16)
    c.rect(94, 658, 12, 12, COLORS["green"], COLORS["green"], 6, 1)
    c.line(100, 664, 100, 805, COLORS["navy"], 4)
    c.arrow(100, 712, 112, 712, COLORS["navy"], 4)
    c.arrow(100, 805, 112, 805, COLORS["navy"], 4)
    c.arrow(434, 712, 516, 738, COLORS["blue"], 4)
    c.arrow(434, 805, 516, 782, COLORS["navy"], 4)
    diamond = c.slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, inch(520), inch(710), inch(152), inch(104))
    diamond.fill.solid(); diamond.fill.fore_color.rgb = rgb(COLORS["soft_amber"])
    diamond.line.color.rgb = rgb(COLORS["amber"]); diamond.line.width = Pt(1.2)
    c.draw.polygon([(596, 710), (672, 762), (596, 814), (520, 762)], fill=hex_tuple(COLORS["soft_amber"]), outline=hex_tuple(COLORS["amber"]))
    c.text(534, 729, 124, 66, "CFG MERGE\nv = v_u + 3(v_c−v_u)", 13, COLORS["ink"], True, "center", "middle")
    c.arrow(676, 762, 716, 762, COLORS["amber"], 4)
    c.node(720, 696, 378, 126, "Joint latent update", "UniPC 同步更新两种生成状态\nvideo latent [1,48,9,33,40]\naction latent [33,64]", COLORS["soft_green"], "90CFAD", COLORS["green"], "Z", 18)

    # Loop is visually separated from the forward path.
    c.path_arrow([(908, 826), (908, 852), (74, 852), (74, 664), (94, 664)], COLORS["green"], 3)
    c.text(456, 832, 344, 20, "STEP 1→4 · 共 8 次 MoT forward", 14, COLORS["green"], True, "center")

    # Step 4: split the final joint state into action and video generation paths.
    c.rect(84, 890, 1044, 108, "F9FBFE", "AFC2D9", 14, 2)
    c.badge(102, 905, 88, "STEP 4", COLORS["amber"])
    c.text(204, 906, 150, 26, "生成输出", 18, COLORS["ink"], True)
    c.path_arrow([(966, 826), (966, 880), (424, 880), (424, 916)], COLORS["green"], 4)
    c.path_arrow([(966, 880), (824, 880), (824, 916)], COLORS["cyan"], 4, True)
    c.rect(358, 920, 360, 60, COLORS["soft_green"], "90CFAD", 11, 2)
    c.text(372, 926, 332, 21, "ACTION GENERATION", 14, COLORS["green"], True)
    c.text(372, 949, 332, 22, "[33,64] → postprocess → robot action [32,8]", 13, COLORS["muted"])
    c.rect(738, 920, 366, 60, COLORS["soft_cyan"], "8CCED9", 11, 2)
    c.text(752, 926, 338, 21, "VIDEO GENERATION（可选）", 14, COLORS["cyan"], True)
    c.text(752, 949, 338, 22, "[1,48,9,33,40] → VAE Decoder → future video", 13, COLORS["muted"])

    # Right comparison panel
    rx, rw = 1208, 656
    c.rect(rx, 166, rw, 422, COLORS["panel"], COLORS["hair"], 18, 2)
    c.badge(rx + 24, 188, 44, "02", COLORS["amber"])
    c.text(rx + 82, 182, 500, 38, "参数与计算结构", 27, COLORS["ink"], True)
    c.text(rx + 24, 229, 606, 24, "采用官方口径；FLOPs 未统一披露，以 denoiser 次数作复杂度 proxy", 15, COLORS["muted"])

    tx, ty = rx + 24, 270
    col = [150, 146, 146, 166]
    headers = ["指标", "Cosmos3 Edge", "Cosmos3 Nano", "LingBot-VA"]
    fills = [COLORS["navy"], COLORS["blue"], COLORS["cyan"], COLORS["green"]]
    xx = tx
    for i, (label, ww) in enumerate(zip(headers, col)):
        c.rect(xx, ty, ww, 44, fills[i], fills[i], 0, 1)
        c.text(xx + 5, ty, ww - 10, 44, label, 15, "FFFFFF", True, "center", "middle")
        xx += ww
    rows = [
        ("总参 / 激活", "4B / dense 2B", "16B / dense 8B", "5.3B\n(5B+~350M)"),
        ("层数 · hidden", "28 · 2048", "36 · 4096", "30 · 3072/768"),
        ("FFN / 双流", "9216 · MoT", "12288 · MoT", "隔离 FFN 双流"),
        ("动作 / 去噪", "32×8 · 4 steps", "32-step profile*", "K=4 · V3/A10"),
        ("请求复杂度", "8 MoT forwards", "8 MoT forwards*", "异步因果双时标"),
    ]
    yy = ty + 44
    for ri, row in enumerate(rows):
        xx = tx
        h = 49
        for ci, (value, ww) in enumerate(zip(row, col)):
            fill = "F6F9FD" if ri % 2 == 0 else "FFFFFF"
            if ci == 0:
                fill = "EDF2F8"
            c.rect(xx, yy, ww, h, fill, COLORS["hair"], 0, 1)
            c.text(xx + 6, yy, ww - 12, h, value, 14 if ci else 15, COLORS["ink"] if ci == 0 else COLORS["muted"], ci == 0, "center", "middle")
            xx += ww
        yy += h
    c.text(tx, 562, 610, 18, "* Nano 为同一 profiling 配置口径；架构与 action chunk 是独立配置。", 13, COLORS["muted"])

    # Right innovation panel
    c.rect(rx, 610, rw, 412, COLORS["panel"], COLORS["hair"], 18, 2)
    c.badge(rx + 24, 632, 44, "03", COLORS["green"])
    c.text(rx + 82, 626, 500, 38, "具身关键创新与竞品定位", 27, COLORS["ink"], True)

    card_y = 680
    cards = [
        ("Action = 状态转移", "以相邻世界状态之间的\n因果 transition 表示 action\n统一视觉与控制时序", COLORS["blue"], COLORS["soft_blue"]),
        ("跨本体几何动作", "ego / effector / grasp\n采用相对 SE(3)；domain-aware\nI/O 保留本体差异", COLORS["amber"], COLORS["soft_amber"]),
        ("FD / ID / Policy 联训", "前向动力学、逆动力学、策略\n共享 MoT；联合预测 future visual\n与 action diffusion", COLORS["green"], COLORS["soft_green"]),
    ]
    x = rx + 24
    for title, detail, accent, fill in cards:
        c.rect(x, card_y, 194, 142, fill, accent, 12, 2)
        c.rect(x, card_y, 194, 7, accent, accent, 3, 1)
        c.text(x + 14, card_y + 18, 166, 28, title, 18, COLORS["ink"], True, "center")
        c.text(x + 12, card_y + 52, 170, 80, detail, 12, COLORS["muted"], False, "center", "middle")
        x += 206

    c.rect(rx + 24, 842, 608, 118, COLORS["soft_navy"], "AFC2D9", 12, 2)
    c.badge(rx + 40, 857, 136, "vs LingBot-VA", COLORS["navy"])
    c.text(rx + 190, 850, 420, 46, "优势在“统一模型迁移闭环”", 20, COLORS["navy"], True, "left", "middle")
    c.text(rx + 40, 900, 564, 48, "Edge：同模型贯通理解→生成→动作；LingBot-VA：面向实时控制的因果双流。\n结论：能力边界更统一 ≠ 已证明实时吞吐更高。", 15, COLORS["muted"], False, "left", "middle")

    # Footer
    c.text(56, 1037, 1808, 22, "Sources: Cosmos 3 paper (arXiv:2606.02800v4) · NVIDIA Edge Policy model card · LingBot-VA paper/repo · local Edge/Nano profiling", 13, COLORS["muted"], False, "left")
    c.text(1720, 1037, 144, 22, "2026.08", 13, COLORS["muted"], True, "right")

    c.set_notes(
        "讲解要点：左侧严格按 STEP 1→4 阅读。先把服务侧预处理折叠成 video/action/prompt 三路输入；"
        "prepare_data_total 分别完成单帧前缀 VAE、文本 tokenize，并按 text→vision→action 顺序打包，同时初始化联合噪声；"
        "sampler_total 的条件与无条件分支读取同一 joint latent，CFG 合流后由 UniPC 更新这个共享状态，再回送给下一步的两条分支；"
        "默认 guidance=3、4 个 UniPC step，因此共 8 次 MoT 前向。最后拆成 action 和可选 video 两路输出。"
        "右侧参数采用 Cosmos 3 与 LingBot-VA 官方论文/模型卡口径。"
        "具身创新聚焦 action-as-transition、跨本体 SE(3) 动作表示，以及 FD/ID/Policy 联合训练。"
    )
    build_h100_slide(c)
    c.save()
    print(PPTX_OUT)
    print(PNG_OUT_PAGE1)
    print(PNG_OUT_PAGE2)


if __name__ == "__main__":
    build()
